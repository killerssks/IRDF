import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import win32com.client

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Use blank slide layout
blank_slide_layout = prs.slide_layouts[6]

# Define Colors
RED_COLOR = RGBColor(214, 18, 18)
BLUE_COLOR = RGBColor(26, 67, 156)
BLACK_COLOR = RGBColor(15, 15, 15)
GREY_COLOR = RGBColor(128, 128, 128)
LIGHT_GREY_COLOR = RGBColor(248, 249, 250)

# Branding paths
logo_cc_path = r"C:\Users\Roshan Kappala\.gemini\antigravity-ide\scratch\logo_careercraft.png"
logo_hcl_path = r"C:\Users\Roshan Kappala\.gemini\antigravity-ide\scratch\logo_hcl.png"
logo_arrow_path = r"C:\Users\Roshan Kappala\.gemini\antigravity-ide\scratch\logo_arrow.png"

def add_base_elements(slide, category):
    # 1. Red circle (large, top right corner)
    red_oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(11.583), Inches(-1.75), Inches(3.5), Inches(3.5)
    )
    red_oval.fill.solid()
    red_oval.fill.fore_color.rgb = RED_COLOR
    red_oval.line.fill.background()
    
    # 2. Blue circle
    blue_oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(11.0), Inches(1.3), Inches(0.6), Inches(0.6)
    )
    blue_oval.fill.solid()
    blue_oval.fill.fore_color.rgb = BLUE_COLOR
    blue_oval.line.fill.background()

    # 3. Subtitle / category at top left
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(10), Inches(0.4))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = category
    p_sub.font.name = "Segoe UI"
    p_sub.font.size = Pt(12)
    p_sub.font.color.rgb = GREY_COLOR

    # 4. Footer images & website link
    if os.path.exists(logo_cc_path):
        slide.shapes.add_picture(logo_cc_path, Inches(0.8), Inches(6.4), width=Inches(2.0))
    if os.path.exists(logo_hcl_path):
        slide.shapes.add_picture(logo_hcl_path, Inches(3.8), Inches(6.45), width=Inches(2.0))
    
    web_box = slide.shapes.add_textbox(Inches(6.8), Inches(6.45), Inches(4.5), Inches(0.4))
    tf_web = web_box.text_frame
    tf_web.word_wrap = True
    tf_web.margin_left = tf_web.margin_top = tf_web.margin_right = tf_web.margin_bottom = 0
    p_web = tf_web.paragraphs[0]
    p_web.text = "www.careercraftacademy.com"
    p_web.font.name = "Segoe UI"
    p_web.font.size = Pt(13)
    p_web.font.bold = True
    p_web.font.color.rgb = BLUE_COLOR
    p_web.alignment = PP_ALIGN.RIGHT

    if os.path.exists(logo_arrow_path):
        slide.shapes.add_picture(logo_arrow_path, Inches(11.8), Inches(6.3), height=Inches(0.6))

def add_slide_title(slide, blue_text, black_text):
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(10), Inches(1.4))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    
    run1 = p_title.add_run()
    run1.text = blue_text
    run1.font.name = "Segoe UI"
    run1.font.size = Pt(36)
    run1.font.bold = True
    run1.font.color.rgb = BLUE_COLOR

    run2 = p_title.add_run()
    run2.text = black_text
    run2.font.name = "Segoe UI"
    run2.font.size = Pt(36)
    run2.font.bold = True
    run2.font.color.rgb = BLACK_COLOR

def add_two_columns_bullets(slide, col1_title, col1_bullets, col2_title, col2_bullets):
    # Left Column
    col1_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(5.6), Inches(3.8))
    tf1 = col1_box.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_top = tf1.margin_right = tf1.margin_bottom = 0
    
    if col1_title:
        p1 = tf1.paragraphs[0]
        p1.text = col1_title
        p1.font.name = "Segoe UI"
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = BLACK_COLOR
        p1.space_after = Pt(12)
        
    for i, bullet in enumerate(col1_bullets):
        if col1_title or i > 0:
            p = tf1.add_paragraph()
        else:
            p = tf1.paragraphs[0]
        p.space_after = Pt(8)
        
        # Add blue checkmark
        r_check = p.add_run()
        r_check.text = "✔  "
        r_check.font.name = "Segoe UI"
        r_check.font.size = Pt(14)
        r_check.font.bold = True
        r_check.font.color.rgb = BLUE_COLOR
        
        r_text = p.add_run()
        r_text.text = bullet
        r_text.font.name = "Segoe UI"
        r_text.font.size = Pt(14)
        r_text.font.color.rgb = BLACK_COLOR
        
    # Right Column
    col2_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.4), Inches(5.6), Inches(3.8))
    tf2 = col2_box.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
    
    if col2_title:
        p2 = tf2.paragraphs[0]
        p2.text = col2_title
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = BLACK_COLOR
        p2.space_after = Pt(12)
        
    for i, bullet in enumerate(col2_bullets):
        if col2_title or i > 0:
            p = tf2.add_paragraph()
        else:
            p = tf2.paragraphs[0]
        p.space_after = Pt(8)
        
        # Add blue checkmark
        r_check = p.add_run()
        r_check.text = "✔  "
        r_check.font.name = "Segoe UI"
        r_check.font.size = Pt(14)
        r_check.font.bold = True
        r_check.font.color.rgb = BLUE_COLOR
        
        r_text = p.add_run()
        r_text.text = bullet
        r_text.font.name = "Segoe UI"
        r_text.font.size = Pt(14)
        r_text.font.color.rgb = BLACK_COLOR

def add_table_slide(slide, headers, rows):
    cols_count = len(headers)
    rows_count = len(rows) + 1
    
    left = Inches(0.8)
    top = Inches(2.4)
    width = Inches(11.7)
    height = Inches(3.8)
    
    table_shape = slide.shapes.add_table(rows_count, cols_count, left, top, width, height)
    table = table_shape.table
    
    # Custom widths
    if cols_count == 3:
        table.columns[0].width = Inches(2.2)
        table.columns[1].width = Inches(4.5)
        table.columns[2].width = Inches(5.0)
    elif cols_count == 2:
        table.columns[0].width = Inches(4.0)
        table.columns[1].width = Inches(7.7)
        
    # Format Headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE_COLOR
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            for r in p.runs:
                r.font.name = "Segoe UI"
                r.font.size = Pt(13)
                r.font.bold = True
                r.font.color.rgb = RGBColor(255, 255, 255)
                
    # Format Data Rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            cell.fill.solid()
            if r_idx % 2 == 0:
                cell.fill.fore_color.rgb = LIGHT_GREY_COLOR
            else:
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name = "Segoe UI"
                    r.font.size = Pt(11)
                    r.font.color.rgb = BLACK_COLOR

def add_grid_cards_slide(slide, cards_data):
    count = len(cards_data)
    if count == 3:
        card_w = Inches(3.6)
        card_h = Inches(3.6)
        spacing = Inches(0.4)
        top = Inches(2.4)
        
        for idx, card in enumerate(cards_data):
            left = Inches(0.8) + idx * (card_w + spacing)
            
            # Rounded rectangle shape
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = LIGHT_GREY_COLOR
            shape.line.color.rgb = RGBColor(220, 224, 230)
            shape.line.width = Pt(1)
            
            tf = shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.3)
            tf.margin_right = Inches(0.3)
            tf.margin_top = Inches(0.3)
            tf.margin_bottom = Inches(0.3)
            
            p_title = tf.paragraphs[0]
            p_title.text = card['title']
            p_title.font.name = "Segoe UI"
            p_title.font.size = Pt(16)
            p_title.font.bold = True
            p_title.font.color.rgb = BLUE_COLOR
            p_title.space_after = Pt(12)
            
            p_text = tf.add_paragraph()
            p_text.text = card['text']
            p_text.font.name = "Segoe UI"
            p_text.font.size = Pt(13)
            p_text.font.color.rgb = BLACK_COLOR

# ----------------- Presentation Deck Configuration -----------------

# Slide 1: Welcome to the Session
s1 = prs.slides.add_slide(blank_slide_layout)
add_base_elements(s1, "Faculty Development Program · 90 Minutes")
add_slide_title(s1, "Demystifying ", "Information Security")
col1_bullets = [
    "Welcome to the session! Congratulations on showing up to learn.",
    "Information security is often seen as exclusively for IT departments, but today we change that.",
    "This presentation deck acts as your personal reference manual for security guidelines."
]
col2_bullets = [
    "Follow along during the session using the sidebar navigations.",
    "Bookmark your links now and share them with colleagues later.",
    "All exercises and labs require only a modern browser — Chrome/Firefox recommended.",
    "Zero software installation is required for any exercises."
]
add_two_columns_bullets(s1, "Welcome to the Session", col1_bullets, "📌 How to Use This Guide", col2_bullets)

# Slide 2: Learning Objectives
s2 = prs.slides.add_slide(blank_slide_layout)
add_base_elements(s2, "Introduction & Setup · Page 1 of 4")
add_slide_title(s2, "Learning ", "Objectives")
col1_bullets = [
    "Explain why universities and academic institutions are high-value targets for cyberattacks.",
    "Define and apply the CIA Triad (Confidentiality, Integrity, Availability) in academia.",
    "Describe the concepts of MFA, Zero Trust, and the Human Firewall.",
    "Check breached credentials using haveibeenpwned.com."
]
col2_bullets = [
    "Operate CyberChef to distinguish encoding, encryption, and hashing.",
    "Analyse raw email headers using free browser tools to identify phishing.",
    "Identify concrete ways to integrate information security topics into course syllabi.",
    "Leave with a personalised, three-step security action plan."
]
add_two_columns_bullets(s2, "Session Outcomes (Part 1)", col1_bullets, "Session Outcomes (Part 2)", col2_bullets)

# Slide 3: Before We Begin — Setup Checklist
s3 = prs.slides.add_slide(blank_slide_layout)
add_base_elements(s3, "Introduction & Setup · Page 1 of 4")
add_slide_title(s3, "Before We Begin ", "Setup Checklist")
col1_bullets = [
    "Open a modern browser. Chrome or Firefox is strongly recommended.",
    "Open Tab 1 — Webbook. Keep the guide open throughout the session.",
    "Open Tab 2 — Have I Been Pwned. Navigate to haveibeenpwned.com.",
    "Open Tab 3 — CyberChef. Navigate to gchq.github.io/CyberChef."
]
col2_bullets = [
    "Open Tab 4 — Google Admin Toolbox. Navigate to toolbox.googleapps.com/apps/messageheader.",
    "Locate a recent email. Find any email in your inbox to copy headers.",
    "Silence notifications but keep Wi-Fi on to avoid distractions.",
    "💡 Pro-Tip: You can use a secondary/university email for labs for privacy."
]
add_two_columns_bullets(s3, "Browser & Active Tabs", col1_bullets, "Data Prep & Notifications", col2_bullets)

# Slide 4: Why Universities Are Prime Targets
s4 = prs.slides.add_slide(blank_slide_layout)
add_base_elements(s4, "Chapter 1 · The Academic Threat Landscape")
add_slide_title(s4, "Why Universities ", "Are Prime Targets")
headers_s4 = ["Factor", "The Academic Reality", "Why Attackers Love It"]
rows_s4 = [
    ["Open culture", "Academic freedom values open sharing and collaboration.", "More entry points, less friction for movement."],
    ["Diverse user base", "Students, faculty, admin, visitors with varying security awareness.", "Weakest link is easier to find and exploit."],
    ["Legacy infrastructure", "Century-old systems alongside bleeding-edge research tools.", "Unpatched legacy systems are vulnerable."],
    ["High-value IP", "Research in pharma, defence, AI, materials science.", "Nation-state actors target research for espionage."],
    ["BYOD culture", "Students and faculty use personal devices on campus Wi-Fi.", "Unmanaged devices introduce vulnerabilities."]
]
add_table_slide(s4, headers_s4, rows_s4)

# Slide 5: Real-World Incidents & Insider Risk
s5 = prs.slides.add_slide(blank_slide_layout)
add_base_elements(s5, "Chapter 1 · The Academic Threat Landscape")
add_slide_title(s5, "Academic Cyber ", "Threats & Incidents")
col1_bullets = [
    "UHI Cyberattack (2021): Closed all 13 campuses and libraries simultaneously.",
    "Indian Universities (2023): Stolen databases of student exams and PII.",
    "Academic attacks increased by 75% globally between 2020 and 2023.",
    "Universities store IP, research, PII, financial systems, and grant data."
]
col2_bullets = [
    "Significant proportion of security incidents result from insider actions.",
    "Typically well-meaning faculty, students, or staff, not malicious intent.",
    "Examples: Clicking phishing links, reusing passwords, misconfiguring cloud storage.",
    "Security awareness is everyone's problem and opportunity to help."
]
add_two_columns_bullets(s5, "Real-World Context", col1_bullets, "The Insider Factor", col2_bullets)

# Slide 6: The CIA Triad — The Foundation of Security
s6 = prs.slides.add_slide(blank_slide_layout)
add_base_elements(s6, "Chapter 1 · The Academic Threat Landscape")
add_slide_title(s6, "The CIA ", "Triad Principle")
cards_s6 = [
    {"title": "🔒 Confidentiality", "text": "Information should only be accessible to those who are authorised to see it. Keep secrets secret. Prevents leaks of personal, research, or financial data."},
    {"title": "✅ Integrity", "text": "Information should be accurate and unaltered. Data should only be changed by authorised people through authorised processes."},
    {"title": "🌐 Availability", "text": "Systems and information should be accessible when authorised users need them. Downtime (e.g. from DDoS or ransomware) is a security failure."}
]
add_grid_cards_slide(s6, cards_s6)

# Slide 7: CIA Triad — The Grading System Analogy
s7 = prs.slides.add_slide(blank_slide_layout)
add_base_elements(s7, "Chapter 1 · The Academic Threat Landscape")
add_slide_title(s7, "Grading System ", "Analogy")
headers_s7 = ["CIA Principle", "Grading System Example", "Breach Scenario"]
rows_s7 = [
    ["🔒 Confidentiality", "Only student, faculty, and controller see grades.", "A grade list leaks on a WhatsApp group before official release."],
    ["✅ Integrity", "A grade of 78 remains 78 — not altered by anyone.", "An attacker hacks the SIS and changes a grade from 45 to 75."],
    ["🌐 Availability", "Exam portal is functional when checking results.", "A DDoS attack crashes the exam portal during results release."]
]
add_table_slide(s7, headers_s7, rows_s7)

# Slide 8: Multi-Factor Authentication (MFA)
s8 = prs.slides.add_slide(blank_slide_layout)
add_base_elements(s8, "Chapter 1 · The Academic Threat Landscape")
add_slide_title(s8, "Multi-Factor ", "Authentication (MFA)")
col1_bullets = [
    "Authentication proves you are who you say you are.",
    "Passwords can be stolen, leaked, or guessed. MFA requires evidence from at least two categories:",
    "🧠 Something you Know: Password, PIN, security question.",
    "📱 Something you Have: OTP, Authenticator app code, security key.",
    "👁️ Something you Are: Fingerprint, face ID, iris scan."
]
col2_bullets = [
    "ATM requires a card (something you have) and a PIN (something you know). One factor is useless without the other.",
    "MFA neutralizes phishing, credential stuffing, and brute force attacks.",
    "Even if an attacker steals your password, they cannot log in without your phone."
]
add_two_columns_bullets(s8, "The Core Concept", col1_bullets, "💡 The ATM Analogy", col2_bullets)

# Slide 9: Zero Trust Architecture & Human Firewall
s9 = prs.slides.add_slide(blank_slide_layout)
add_base_elements(s9, "Chapter 1 · The Academic Threat Landscape")
add_slide_title(s9, "Zero Trust & ", "The Human Firewall")
col1_bullets = [
    "Core Principle: 'Never Trust, Always Verify'.",
    "No user/device is trusted by default, regardless of network location.",
    "Every access request is verified every single time.",
    "Access is granted under 'least privilege' (minimum needed resources).",
    "Checks device compliance, location, timing, and anomalies."
]
col2_bullets = [
    "A firewall monitors network traffic; a human firewall is security-aware people.",
    "Verizon DBIR: 74% to 82% of all data breaches involve a human element.",
    "Technical controls are useless if a user clicks a phishing link or shares passwords.",
    "Humans are the greatest vulnerability AND the most potent defence."
]
add_two_columns_bullets(s9, "Zero Trust Architecture", col1_bullets, "The Human Firewall", col2_bullets)

# Slide 10: Lab 1: The Credential Check (HIBP)
s10 = prs.slides.add_slide(blank_slide_layout)
add_base_elements(s10, "Chapter 2 · Defender's Web Toolkit")
add_slide_title(s10, "Lab 1: Have I ", "Been Pwned")
col1_bullets = [
    "Data Breach: Stolen user credentials sold or published on the dark web.",
    "Credential Stuffing: Stolen credentials from Site A are used to attack Site B.",
    "If you reuse passwords, a breach at a minor site compromises your university portal.",
    "Database: HIBP holds over 13 billion compromised accounts."
]
col2_bullets = [
    "1. Navigate to haveibeenpwned.com.",
    "2. Type your primary email in the search bar.",
    "3. Click the 'pwned?' button.",
    "4. Review the details: name of breached service, date, and data type compromised.",
    "5. Search other personal and work email addresses."
]
add_two_columns_bullets(s10, "Data Breaches & Password Reuse", col1_bullets, "Step-by-Step Exercise", col2_bullets)

# Slide 11: If You've Been Pwned — Action Plan
s11 = prs.slides.add_slide(blank_slide_layout)
add_base_elements(s11, "Chapter 2 · Defender's Web Toolkit")
add_slide_title(s11, "Breached Account ", "Action Plan")
col1_bullets = [
    "Change your password on each breached service immediately.",
    "Identify which accounts used that password and change them too.",
    "Focus on breaches that exposed passwords first.",
    "If the service no longer exists, just note the password it used."
]
col2_bullets = [
    "Enable Multi-Factor Authentication (MFA) on your critical accounts immediately.",
    "Adopt a password manager to handle unique, complex passwords.",
    "Bitwarden is free and open-source; 1Password is an excellent alternative.",
    "Ask your IT department about institutional licensing."
]
add_two_columns_bullets(s11, "1. Change and Audit Passwords", col1_bullets, "2. Secure and Manage Accounts", col2_bullets)

# Slide 12: Lab 2: CyberChef Introduction
s12 = prs.slides.add_slide(blank_slide_layout)
add_base_elements(s12, "Chapter 2 · Defender's Web Toolkit")
add_slide_title(s12, "Lab 2: CyberChef ", "Intro")
col1_bullets = [
    "Created and open-sourced by GCHQ (UK Intelligence agency).",
    "Intended for carrying out 'cyber' operations in a web browser.",
    "Extremely powerful for data transforming, encoding, and encrypting.",
    "Runs entirely locally — no data leaves your machine.",
    "Uses a 'recipe' metaphor: drag operations, input data, see output."
]
col2_bullets = [
    "Encoding (Compatibility): Transforms format, not for security (e.g. Morse code, Base64). Easily reversed by anyone.",
    "Encryption (Confidentiality): Unreadable without correct key (e.g. AES). Computationally infeasible to read without key.",
    "Hashing (Integrity): One-way fingerprint of data. Used to verify integrity, not for secrecy (e.g. SHA-256)."
]
add_two_columns_bullets(s12, "What is CyberChef?", col1_bullets, "Encoding vs. Encryption vs. Hashing", col2_bullets)

# Slide 13: CyberChef Exercises (Base64, AES, Hashing)
s13 = prs.slides.add_slide(blank_slide_layout)
add_base_elements(s13, "Chapter 2 · Defender's Web Toolkit")
add_slide_title(s13, "CyberChef in ", "Practice")
col1_bullets = [
    "Exercise 2A (Base64): Decode SGVsbG8gV29ybGQ=. Drag 'From Base64'. Output is 'Hello World'. Confirms Base64 is not secure.",
    "Exercise 2B (AES): Input sensitive grade data. Use Key: MySecretKey12345 (UTF8, CBC). Observe ciphertext output.",
    "Test decryption: Switch to 'AES Decrypt'. Use same key. Decrypts back to grade data."
]
col2_bullets = [
    "Exercise 2C (Hashing): Input 'The exam is on Tuesday at 10am' into MD5.",
    "Observe 32-char hex hash output.",
    "Avalanche Effect: Change 'T' to 't'. Notice the hash changes completely.",
    "SHA-256 Hashing: Replace MD5 with SHA-256. Notice the 64-char stronger hash.",
    "Integrity verification: Compare downloaded file hashes to official hashes."
]
add_two_columns_bullets(s13, "Base64 & AES Exercises", col1_bullets, "Hashing & Integrity Exercise", col2_bullets)

# Slide 14: Lab 3: X-Raying Phishing Emails
s14 = prs.slides.add_slide(blank_slide_layout)
add_base_elements(s14, "Chapter 2 · Defender's Web Toolkit")
add_slide_title(s14, "Lab 3: Email ", "Headers Analysis")
col1_bullets = [
    "Email headers are invisible metadata documenting the journey of the email.",
    "Shows every mail server passed, timestamps, and origin IP.",
    "Email clients hide headers to keep the interface clean.",
    "Attackers rely on this invisibility to spoof senders.",
    "Headers are read bottom-up: top headers are added by trusted servers."
]
col2_bullets = [
    "Gmail: Open email → Click 3-dot menu → Select 'Show original' → Copy.",
    "Outlook Web: Open email → Click 3-dot menu → View → 'View message source'.",
    "University Webmail (Roundcube/Zimbra): Search for 'View Source', 'Show Raw', or 'All Headers'.",
    "Copy the entire raw block of text to clipboard."
]
add_two_columns_bullets(s14, "Understanding Email Headers", col1_bullets, "Extracting Raw Headers", col2_bullets)

# Slide 15: Google Admin Toolbox & Phishing Checklist
s15 = prs.slides.add_slide(blank_slide_layout)
add_base_elements(s15, "Chapter 2 · Defender's Web Toolkit")
add_slide_title(s15, "Phishing Header ", "Analysis")
headers_s15 = ["Header Field", "What It Tells You", "Can It Be Faked?"]
rows_s15 = [
    ["From:", "The sender's display name and address", "✔ Yes — easily spoofed"],
    ["Reply-To:", "Where replies will go (often different in phishing)", "✔ Yes"],
    ["Received:", "Mail servers chain of custody", "⚠ Lower entries can be faked; top cannot"],
    ["Return-Path:", "Delivery failure address — true sender domain", "⚠ Often inconsistent in phishing"],
    ["SPF/DKIM/DMARC", "Authentication results — verified domain", "❌ No — verified by server"],
    ["Message-ID:", "Unique identifier — check if domain matches From", "⚠ Often inconsistent in phishing"]
]
add_table_slide(s15, headers_s15, rows_s15)

# Slide 16: Bringing It to the Classroom (Business & Law)
s16 = prs.slides.add_slide(blank_slide_layout)
add_base_elements(s16, "Chapter 3 · Classroom Integration")
add_slide_title(s16, "Classroom ", "Integration (Part 1)")
col1_bullets = [
    "Corporate Governance: Case studies on Sony hack (2014) or Air India breach (2021) for board-level liability.",
    "Marketing: Audit company privacy policies and consent flows.",
    "Supply Chain: Vendor risk analysis (SolarWinds attack).",
    "Entrepreneurship: Add 'Security & Privacy' sections to business plans."
]
col2_bullets = [
    "Evidence: Discuss digital evidence admissibility, chain of custody, and hashing.",
    "Privacy Law: Mock data breach response and notification letters under DPDP Act 2023.",
    "Cyberlaw: Apply IT Act 2000 (Sec 43, 66) to credential stuffing cases."
]
add_two_columns_bullets(s16, "Business & Management", col1_bullets, "Law & Legal Studies", col2_bullets)

# Slide 17: Bringing It to the Classroom (Healthcare & Humanities)
s17 = prs.slides.add_slide(blank_slide_layout)
add_base_elements(s17, "Chapter 3 · Classroom Integration")
add_slide_title(s17, "Classroom ", "Integration (Part 2)")
col1_bullets = [
    "Medical Ethics: Ransomware attacks affecting patient record availability.",
    "Public Health: Data integrity in outbreak surveillance systems.",
    "Clinical Research: IRB data security compliance and encrypted sharing."
]
col2_bullets = [
    "Media Studies: Expose metadata (EXIF) to check image manipulation.",
    "Sociology: Geopolitics of encryption and state-citizen power balances.",
    "Education: Check student privacy in EdTech tools (what data is stored?)."
]
add_two_columns_bullets(s17, "Healthcare & Life Sciences", col1_bullets, "Humanities & Social Sciences", col2_bullets)

# Slide 18: 3-Step Personal Security Action Plan
s18 = prs.slides.add_slide(blank_slide_layout)
add_base_elements(s18, "Chapter 3 · Classroom Integration")
add_slide_title(s18, "Personal Security ", "Action Plan")
cards_s18 = [
    {"title": "Step 1: Secure Accounts", "text": "Enable MFA on university and personal emails. Use authenticator apps over SMS. Change passwords flagged by haveibeenpwned.com immediately."},
    {"title": "Step 2: Password Manager", "text": "Adopt Bitwarden (free/open-source) to store and generate unique, strong passwords. Stop reusing passwords across platforms."},
    {"title": "Step 3: Teach One Colleague", "text": "Share this webbook and help one colleague check Have I Been Pwned. A stronger human firewall protects the whole institution."}
]
add_grid_cards_slide(s18, cards_s18)

# Slide 19: Further Learning & Resources
s19 = prs.slides.add_slide(blank_slide_layout)
add_base_elements(s19, "Chapter 3 · Classroom Integration")
add_slide_title(s19, "Resources & ", "Further Learning")
col1_bullets = [
    "haveibeenpwned.com — Credential breach checker.",
    "gchq.github.io/CyberChef — Hashing and encryption.",
    "toolbox.googleapps.com — Email header analyser.",
    "Bitwarden — Open-source password manager.",
    "Coursera Google Cybersecurity / Cybrary / SANS Cyber Aces."
]
col2_bullets = [
    "CERT-In — India's cyber incident response team.",
    "DPDP Act 2023 — Digital Personal Data Protection Act India.",
    "The Art of Invisibility by Kevin Mitnick.",
    "Countdown to Zero Day by Kim Zetter (Stuxnet story).",
    "Sandworm by Andy Greenberg (Nation-state threats).",
    "This Is How They Tell Me the World Ends by Nicole Perlroth."
]
add_two_columns_bullets(s19, "Tools & Certifications", col1_bullets, "Policy & Recommended Books", col2_bullets)

# Save the presentation
output_pptx_path = r"c:\Users\Roshan Kappala\OneDrive\Desktop\Rohith\Demystifying_Information_Security.pptx"
prs.save(output_pptx_path)
print(f"Successfully generated 19 slides in PPTX format: {output_pptx_path}")
